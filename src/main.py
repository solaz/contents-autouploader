"""Main CLI interface and pipeline orchestrator."""

import json
from pathlib import Path

import click
from rich.console import Console
from rich.progress import Progress, SpinnerColumn, TextColumn
from rich.panel import Panel

from src.config import init_settings, settings
from src.generators.ppt_generator import PPTGenerator
from src.generators.script_generator import ScriptGenerator
from src.generators.tts_generator import TTSGenerator
from src.generators.video_generator import VideoGenerator
from src.models.script import Script, ScriptInput
from src.services.sync_service import SyncService
from src.services.youtube_service import YouTubeService
from src.utils.helpers import ensure_dir, sanitize_filename

console = Console()


@click.group()
@click.option("--config", "-c", type=click.Path(exists=True), help="Path to config file")
@click.pass_context
def cli(ctx, config):
    """Contents Auto-Uploader: Automated content creation pipeline."""
    ctx.ensure_object(dict)
    if config:
        init_settings(config)
    ctx.obj["settings"] = settings()


@cli.command()
@click.option("--topic", "-t", required=True, help="Main topic of the presentation")
@click.option("--storyline", "-s", required=True, help="Brief storyline or outline")
@click.option("--duration", "-d", default=10, help="Target duration in minutes")
@click.option("--tone", default="친근하고 설득력 있는", help="Tone and style")
@click.option("--output", "-o", type=click.Path(), help="Output JSON file path")
@click.pass_context
def script(ctx, topic, storyline, duration, tone, output):
    """Generate a script from topic and storyline."""
    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        console=console,
    ) as progress:
        task = progress.add_task("Generating script...", total=None)

        generator = ScriptGenerator(config=ctx.obj["settings"])
        script_input = ScriptInput(
            topic=topic,
            storyline=storyline,
            duration_minutes=duration,
            tone=tone,
        )

        result = generator.generate(script_input)
        progress.update(task, completed=True)

    # Output
    if output:
        output_path = Path(output)
    else:
        output_dir = ensure_dir(Path(settings().output.base_dir) / "scripts")
        output_path = output_dir / f"{sanitize_filename(topic)}.json"

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result.model_dump(), f, ensure_ascii=False, indent=2)

    console.print(Panel(f"Script saved to: {output_path}", title="Success", style="green"))
    console.print(f"Total sections: {len(result.sections)}")
    console.print(f"Estimated duration: {result.total_duration_sec:.1f} seconds")


@cli.command()
@click.option("--script", "-s", required=True, type=click.Path(exists=True), help="Path to script JSON")
@click.option("--output", "-o", type=click.Path(), help="Output PPTX file path")
@click.pass_context
def ppt(ctx, script, output):
    """Generate PowerPoint presentation from script."""
    with open(script, encoding="utf-8") as f:
        script_data = json.load(f)

    script_obj = Script(**script_data)

    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        console=console,
    ) as progress:
        task = progress.add_task("Generating presentation...", total=None)

        generator = PPTGenerator(config=ctx.obj["settings"])
        presentation = generator.generate(script_obj, output)
        progress.update(task, completed=True)

    console.print(Panel(f"Presentation saved to: {presentation.file_path}", title="Success", style="green"))
    console.print(f"Total slides: {len(presentation.slides)}")


@cli.command()
@click.option("--script", "-s", required=True, type=click.Path(exists=True), help="Path to script JSON")
@click.option("--output-dir", "-o", type=click.Path(), help="Output directory for audio files")
@click.option("--provider", "-p", type=click.Choice(["elevenlabs", "google", "openai"]), help="TTS provider")
@click.pass_context
def tts(ctx, script, output_dir, provider):
    """Generate TTS audio from script."""
    with open(script, encoding="utf-8") as f:
        script_data = json.load(f)

    script_obj = Script(**script_data)

    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        console=console,
    ) as progress:
        task = progress.add_task("Generating audio...", total=None)

        generator = TTSGenerator(config=ctx.obj["settings"], provider=provider)
        results = generator.generate_for_script(script_obj, output_dir)
        progress.update(task, completed=True)

    total_duration = sum(duration for _, _, duration in results)
    console.print(Panel(f"Audio files generated: {len(results)}", title="Success", style="green"))
    console.print(f"Total duration: {total_duration:.1f} seconds")

    for section_id, audio_path, duration in results:
        console.print(f"  Section {section_id}: {audio_path.name} ({duration:.1f}s)")


@cli.command()
@click.option("--script", "-s", required=True, type=click.Path(exists=True), help="Path to script JSON")
@click.option("--ppt", "-p", type=click.Path(exists=True), help="Path to PPTX file (optional, will generate if not provided)")
@click.option("--audio-dir", "-a", type=click.Path(exists=True), help="Directory with audio files (optional, will generate if not provided)")
@click.option("--output", "-o", type=click.Path(), help="Output video file path")
@click.option("--transitions/--no-transitions", default=True, help="Enable fade transitions")
@click.pass_context
def video(ctx, script, ppt, audio_dir, output, transitions):
    """Generate video from script, presentation, and audio."""
    config = ctx.obj["settings"]

    with open(script, encoding="utf-8") as f:
        script_data = json.load(f)
    script_obj = Script(**script_data)

    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        console=console,
    ) as progress:
        # Generate or load presentation
        if ppt:
            from pptx import Presentation as PPTXPresentation
            from src.models.presentation import Presentation, Slide

            pptx = PPTXPresentation(ppt)
            slides = []
            for i, slide in enumerate(pptx.slides):
                title = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        title = shape.text_frame.text
                        break
                slides.append(Slide(slide_index=i, title=title, content=[], notes=""))
            presentation = Presentation(title=script_obj.title, slides=slides, file_path=Path(ppt))
        else:
            task = progress.add_task("Generating presentation...", total=None)
            ppt_generator = PPTGenerator(config=config)
            presentation = ppt_generator.generate(script_obj)
            progress.update(task, completed=True)

        # Export slides as images
        task = progress.add_task("Exporting slides as images...", total=None)
        ppt_generator = PPTGenerator(config=config)
        image_paths = ppt_generator.export_slides_as_images(presentation)
        progress.update(task, completed=True)

        # Generate or load audio
        if audio_dir:
            audio_dir = Path(audio_dir)
            audio_results = []
            
            # Need a way to measure duration if loading existing files
            try:
                from mutagen.mp3 import MP3
                from mutagen.wave import WAVE
            except ImportError:
                console.print("[yellow]Warning: mutagen not installed. Audio duration estimation might be inaccurate.[/yellow]")

            for section in script_obj.sections:
                # Try mp3 then wav
                audio_path_mp3 = audio_dir / f"section_{section.section_id:03d}.mp3"
                audio_path_wav = audio_dir / f"section_{section.section_id:03d}.wav"
                
                audio_path = None
                duration = section.estimated_duration_sec  # Default fallback
                
                if audio_path_mp3.exists():
                    audio_path = audio_path_mp3
                    try:
                        duration = MP3(str(audio_path)).info.length
                    except:
                        pass
                elif audio_path_wav.exists():
                    audio_path = audio_path_wav
                    try:
                        duration = WAVE(str(audio_path)).info.length
                    except:
                        pass
                
                if audio_path:
                    audio_results.append((section.section_id, audio_path, duration))
        else:
            task = progress.add_task("Generating audio...", total=None)
            tts_generator = TTSGenerator(config=config)
            audio_results = tts_generator.generate_for_script(script_obj)
            progress.update(task, completed=True)

        # Create sync data
        task = progress.add_task("Synchronizing...", total=None)
        sync_service = SyncService()
        sync_data = sync_service.create_sync_data(script_obj, presentation, audio_results)
        progress.update(task, completed=True)

        # Generate video
        task = progress.add_task("Generating video...", total=None)
        video_generator = VideoGenerator(config=config)
        if transitions:
            video_path = video_generator.generate_with_transitions(presentation, sync_data, output)
        else:
            video_path = video_generator.generate(presentation, sync_data, output)
        progress.update(task, completed=True)

    console.print(Panel(f"Video saved to: {video_path}", title="Success", style="green"))
    console.print(f"Duration: {sync_data.total_duration:.1f} seconds")


@cli.command()
@click.option("--video", "-v", required=True, type=click.Path(exists=True), help="Path to video file")
@click.option("--title", "-t", required=True, help="Video title")
@click.option("--description", "-d", help="Video description")
@click.option("--tags", help="Comma-separated tags")
@click.option("--privacy", type=click.Choice(["public", "private", "unlisted"]), default="private")
@click.option("--thumbnail", type=click.Path(exists=True), help="Thumbnail image path")
@click.pass_context
def upload(ctx, video, title, description, tags, privacy, thumbnail):
    """Upload video to YouTube."""
    tags_list = [t.strip() for t in tags.split(",")] if tags else None

    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        console=console,
    ) as progress:
        task = progress.add_task("Uploading to YouTube...", total=None)

        service = YouTubeService(config=ctx.obj["settings"])
        video_id = service.upload(
            video_path=video,
            title=title,
            description=description or "",
            tags=tags_list,
            privacy_status=privacy,
            thumbnail_path=thumbnail,
        )
        progress.update(task, completed=True)

    video_url = service.get_video_url(video_id)
    console.print(Panel(f"Video uploaded!\nURL: {video_url}", title="Success", style="green"))


@cli.command()
@click.option("--topic", "-t", required=True, help="Main topic of the presentation")
@click.option("--storyline", "-s", required=True, help="Brief storyline or outline")
@click.option("--duration", "-d", default=10, help="Target duration in minutes")
@click.option("--tone", default="친근하고 설득력 있는", help="Tone and style")
@click.option("--upload/--no-upload", default=False, help="Upload to YouTube after generation")
@click.option("--privacy", type=click.Choice(["public", "private", "unlisted"]), default="private")
@click.pass_context
def generate(ctx, topic, storyline, duration, tone, upload, privacy):
    """Run the complete content generation pipeline."""
    config = ctx.obj["settings"]

    console.print(Panel(f"Starting content generation pipeline\nTopic: {topic}", title="Pipeline", style="blue"))

    with Progress(
        SpinnerColumn(),
        TextColumn("[progress.description]{task.description}"),
        console=console,
    ) as progress:
        # Step 1: Generate script
        task = progress.add_task("Step 1/5: Generating script...", total=None)
        script_generator = ScriptGenerator(config=config)
        script_input = ScriptInput(
            topic=topic,
            storyline=storyline,
            duration_minutes=duration,
            tone=tone,
        )
        script_obj = script_generator.generate(script_input)
        progress.update(task, completed=True)
        console.print(f"  Generated {len(script_obj.sections)} sections")

        # Step 2: Generate presentation
        task = progress.add_task("Step 2/5: Generating presentation...", total=None)
        ppt_generator = PPTGenerator(config=config)
        presentation = ppt_generator.generate(script_obj)
        progress.update(task, completed=True)
        console.print(f"  Created {len(presentation.slides)} slides")

        # Step 3: Export slides as images
        task = progress.add_task("Step 3/5: Exporting slide images...", total=None)
        image_paths = ppt_generator.export_slides_as_images(presentation)
        progress.update(task, completed=True)

        # Step 4: Generate TTS audio
        task = progress.add_task("Step 4/5: Generating audio...", total=None)
        tts_generator = TTSGenerator(config=config)
        audio_results = tts_generator.generate_for_script(script_obj)
        progress.update(task, completed=True)
        total_audio_duration = sum(d for _, _, d in audio_results)
        console.print(f"  Total audio duration: {total_audio_duration:.1f}s")

        # Create sync data
        sync_service = SyncService()
        sync_data = sync_service.create_sync_data(script_obj, presentation, audio_results)

        # Step 5: Generate video
        task = progress.add_task("Step 5/5: Generating video...", total=None)
        video_generator = VideoGenerator(config=config)
        video_path = video_generator.generate_with_transitions(presentation, sync_data)
        progress.update(task, completed=True)

    console.print(Panel(f"Video generated: {video_path}", title="Complete", style="green"))

    # Optional: Upload to YouTube
    if upload:
        console.print("\n[bold]Uploading to YouTube...[/bold]")
        service = YouTubeService(config=config)
        video_id = service.upload(
            video_path=video_path,
            title=script_obj.title,
            description=script_obj.description,
            tags=script_obj.tags,
            privacy_status=privacy,
        )
        video_url = service.get_video_url(video_id)
        console.print(Panel(f"Uploaded to YouTube!\nURL: {video_url}", title="YouTube", style="green"))


if __name__ == "__main__":
    cli()
