"""PowerPoint presentation generator."""

from pathlib import Path

from pptx import Presentation as PPTXPresentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from src.config import Settings, settings
from src.models.presentation import Presentation, Slide
from src.models.script import Script
from src.utils.helpers import ensure_dir, sanitize_filename


class PPTGenerator:
    """Generator for creating PowerPoint presentations from scripts."""

    def __init__(self, config: Settings | None = None):
        self.config = config or settings()
        self.ppt_config = self.config.presentation

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGB."""
        hex_color = hex_color.lstrip("#")
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    def generate(self, script: Script, output_path: Path | str | None = None) -> Presentation:
        """Generate a PowerPoint presentation from a script."""
        prs = PPTXPresentation()

        # Set slide dimensions (16:9)
        prs.slide_width = Inches(self.ppt_config.width / 96)  # Convert pixels to inches
        prs.slide_height = Inches(self.ppt_config.height / 96)

        slides = []

        # Create title slide
        title_slide = self._create_title_slide(prs, script.title, script.description)
        slides.append(
            Slide(
                slide_index=0,
                title=script.title,
                content=[script.description] if script.description else [],
                notes=f"안녕하세요, 오늘은 {script.title}에 대해 이야기해보겠습니다.",
            )
        )

        # Create content slides from sections
        for i, section in enumerate(script.sections):
            content_slide = self._create_content_slide(
                prs,
                section.title,
                section.key_points,
                section.content,
            )
            slides.append(
                Slide(
                    slide_index=i + 1,
                    title=section.title,
                    content=section.key_points,
                    notes=section.content,
                )
            )

        # Determine output path
        if output_path is None:
            output_dir = ensure_dir(Path(self.config.output.base_dir) / "presentations")
            filename = sanitize_filename(script.title) + ".pptx"
            output_path = output_dir / filename
        else:
            output_path = Path(output_path)
            ensure_dir(output_path.parent)

        # Save presentation
        prs.save(str(output_path))

        return Presentation(
            title=script.title,
            slides=slides,
            file_path=output_path,
        )

    def _create_title_slide(
        self,
        prs: PPTXPresentation,
        title: str,
        subtitle: str,
    ):
        """Create the title slide."""
        # Use blank layout and add shapes manually for more control
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(3),
            prs.slide_width - Inches(1),
            Inches(1.5),
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(self.ppt_config.title_font_size + 10)
        title_para.font.bold = True
        title_para.font.color.rgb = self._hex_to_rgb(self.ppt_config.title_color)
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1),
                Inches(4.5),
                prs.slide_width - Inches(2),
                Inches(1),
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(self.ppt_config.body_font_size)
            subtitle_para.font.color.rgb = self._hex_to_rgb(self.ppt_config.body_color)
            subtitle_para.alignment = PP_ALIGN.CENTER

        return slide

    def _create_content_slide(
        self,
        prs: PPTXPresentation,
        title: str,
        bullet_points: list[str],
        notes: str,
    ):
        """Create a content slide with title and bullet points."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(0.5),
            prs.slide_width - Inches(1),
            Inches(1),
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(self.ppt_config.title_font_size)
        title_para.font.bold = True
        title_para.font.color.rgb = self._hex_to_rgb(self.ppt_config.title_color)

        # Add bullet points
        if bullet_points:
            content_box = slide.shapes.add_textbox(
                Inches(0.75),
                Inches(1.75),
                prs.slide_width - Inches(1.5),
                Inches(5),
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for i, point in enumerate(bullet_points):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()

                para.text = f"• {point}"
                para.font.size = Pt(self.ppt_config.body_font_size)
                para.font.color.rgb = self._hex_to_rgb(self.ppt_config.body_color)
                para.space_after = Pt(12)

        # Add speaker notes
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

        return slide

    def export_slides_as_images(
        self,
        presentation: Presentation,
        output_dir: Path | str | None = None,
    ) -> list[Path]:
        """Export slides as images for video generation.

        Note: This requires LibreOffice or unoconv to be installed.
        Alternative: Use python-pptx-to-image or similar library.
        """
        if presentation.file_path is None:
            raise ValueError("Presentation file path is not set")

        if output_dir is None:
            output_dir = ensure_dir(
                Path(self.config.output.base_dir) / "slides" / sanitize_filename(presentation.title)
            )
        else:
            output_dir = ensure_dir(Path(output_dir))

        # For now, we'll use a placeholder that creates simple images
        # In production, you'd use LibreOffice or a conversion service
        image_paths = []

        try:
            from PIL import Image, ImageDraw, ImageFont
        except ImportError:
            raise ImportError("Pillow is required for image generation")

        for slide in presentation.slides:
            img = Image.new(
                "RGB",
                (self.ppt_config.width, self.ppt_config.height),
                color=self.ppt_config.background_color,
            )
            draw = ImageDraw.Draw(img)

            # Try to use a system font, fall back to default
            try:
                title_font = ImageFont.truetype("/System/Library/Fonts/AppleSDGothicNeo.ttc", 60)
                body_font = ImageFont.truetype("/System/Library/Fonts/AppleSDGothicNeo.ttc", 36)
            except (OSError, IOError):
                title_font = ImageFont.load_default()
                body_font = ImageFont.load_default()

            # Draw title
            title_color = tuple(
                int(self.ppt_config.title_color.lstrip("#")[i : i + 2], 16) for i in (0, 2, 4)
            )
            draw.text((100, 80), slide.title, font=title_font, fill=title_color)

            # Draw content
            body_color = tuple(
                int(self.ppt_config.body_color.lstrip("#")[i : i + 2], 16) for i in (0, 2, 4)
            )
            y_pos = 200
            for point in slide.content:
                draw.text((120, y_pos), f"• {point}", font=body_font, fill=body_color)
                y_pos += 60

            # Save image
            image_path = output_dir / f"slide_{slide.slide_index:03d}.png"
            img.save(str(image_path))
            image_paths.append(image_path)

            # Update slide with image path
            slide.image_path = image_path

        return image_paths
